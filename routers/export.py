from __future__ import annotations

from fastapi import APIRouter
from fastapi.responses import StreamingResponse, JSONResponse
from pydantic import BaseModel
from io import BytesIO
from pptx import Presentation
from pptx.util import Pt, Inches
from pptx.dml.color import RGBColor
import os
from collections import Counter
from typing import List, Dict, Any
import httpx
import re

router = APIRouter(prefix="/export", tags=["export"])

# ─────────────────────────────────────────────────────────
# 0) 괄호 감지/정리/전각치환 유틸
# ─────────────────────────────────────────────────────────

# ASCII ()[]{} + FULLWIDTH （）［］｛｝
_BR_ANY = re.compile(r"[\(\)\[\]\{\}\uFF08\uFF09\uFF3B\uFF3D\uFF5B\uFF5D]")

def has_any_bracket(s: str) -> bool:
    return bool(_BR_ANY.search(s or ""))

def normalize_bracket_spacing(s: str) -> str:
    """괄호 안/앞뒤 불필요한 공백과 문장부호 앞 공백 정리."""
    if not s:
        return ""
    s = re.sub(r'([\(\[\{])\s+', r'\1', s)        # 여는 괄호 뒤 공백 제거
    s = re.sub(r'\s+([\)\]\}])', r'\1', s)        # 닫는 괄호 앞 공백 제거
    s = re.sub(r'\s+([,.;:!?])', r'\1', s)        # 문장부호 앞 공백 제거
    s = re.sub(r'[ \t\u00A0]+', ' ', s)           # 연속 공백을 하나로
    s = re.sub(r' {2,}', ' ', s).strip()
    return s

# PPT 렌더러에서 괄호 내 텍스트 누락 방지용(전각 치환)
USE_FULLWIDTH_BRACKETS = True
_BR_MAP = str.maketrans({
    "(": "（", ")": "）",
    "[": "［", "]": "］",
    "{": "｛", "}": "｝",
})
def safe_brackets(s: str) -> str:
    if not USE_FULLWIDTH_BRACKETS:
        return s or ""
    return (s or "").translate(_BR_MAP)

# ─────────────────────────────────────────────────────────
# 0-1) 괄호별 색상 유틸 (한 군데에서만 정의)
# ─────────────────────────────────────────────────────────
_COLORS = {
    "square":  RGBColor(25, 118, 210),  # [] 파랑  #1976D2
    "round":   RGBColor(0, 138, 0),     # () 초록  #008A00
    "curly":   RGBColor(230, 81, 0),    # {} 주황  #E65100
    "default": RGBColor(0, 0, 0),       # 일반 텍스트(검정)
}
_BLACK = RGBColor(0, 0, 0)

def colorize_brackets(paragraph, text: str, size_pt: int = 20):
    """
    괄호 문자만 컬러, 나머지 모든 텍스트는 검정으로 출력.
    전각/반각 괄호 모두 처리.
    """
    paragraph.text = ""

    opens = {"[": "square", "(": "round", "{": "curly",
             "［": "square", "（": "round", "｛": "curly"}
    closes = {"]": "square", ")": "round", "}": "curly",
              "］": "square", "）": "round", "｝": "curly"}

    buf = []
    def flush_buf():
        if buf:
            run = paragraph.add_run()
            run.text = "".join(buf)
            run.font.size = Pt(size_pt)
            run.font.color.rgb = _BLACK
            buf.clear()

    for ch in (text or ""):
        if ch in opens:
            flush_buf()
            run = paragraph.add_run()
            run.text = ch
            run.font.size = Pt(size_pt)
            run.font.color.rgb = _COLORS[opens[ch]]
        elif ch in closes:
            flush_buf()
            run = paragraph.add_run()
            run.text = ch
            run.font.size = Pt(size_pt)
            run.font.color.rgb = _COLORS[closes[ch]]
        else:
            buf.append(ch)

    flush_buf()

# (옵션) 블록 단위 색칠 버전 – 현재 사용하지 않지만 남겨둠
_BR_SPAN = re.compile(r'(\[[^\]]*\]|［[^］]*］|\([^\)]*\)|（[^）]*）|\{[^}]*\}|｛[^｝]*｝)')
def add_colored_runs(paragraph, text: str, size_pt: int = 20):
    if hasattr(paragraph, "clear"):
        paragraph.clear()
    else:
        paragraph.text = ""

    pos = 0
    s = text or ""
    for m in _BR_SPAN.finditer(s):
        if m.start() > pos:
            r = paragraph.add_run()
            r.text = s[pos:m.start()]
            r.font.color.rgb = _COLORS["default"]
            r.font.size = Pt(size_pt)

        seg = m.group(0)
        r = paragraph.add_run()
        r.text = seg
        if seg.startswith(("［", "[")):
            r.font.color.rgb = _COLORS["square"]
        elif seg.startswith(("（", "(")):
            r.font.color.rgb = _COLORS["round"]
        elif seg.startswith(("｛", "{")):
            r.font.color.rgb = _COLORS["curly"]
        else:
            r.font.color.rgb = _COLORS["default"]
        r.font.size = Pt(size_pt)
        pos = m.end()

    if pos < len(s):
        r = paragraph.add_run()
        r.text = s[pos:]
        r.font.color.rgb = _COLORS["default"]
        r.font.size = Pt(size_pt)

# ─────────────────────────────────────────────────────────
# 1) (데모) 분석 스텁 — 실제 서비스 함수로 교체 가능
# ─────────────────────────────────────────────────────────
def analyze_paragraph(text: str) -> dict:
    n = max(len(text), 3)
    a = int(n * 0.33); b = int(n * 0.66)
    return {"outline": {"intro": text[:a].strip(),
                        "body":  text[a:b].strip(),
                        "conclusion": text[b:].strip()}}

def analyze_topic_title_summary(text: str) -> dict:
    return {
        "topic_en": "Secret marketing tactics in retail",
        "topic_ko": "소매 유통의 비밀 마케팅 전술",
        "title_en": "The Allure of Secret Menus and Sales",
        "title_ko": "비밀 메뉴와 세일의 매혹",
        "gist_en": ("Secret menus and early sale access create exclusivity, "
                    "enhancing customer loyalty and excitement in retail."),
        "gist_ko": "비밀 메뉴와 조기 세일 접근은 배타성을 만들어 고객 충성도와 흥미를 높입니다.",
    }

# ─────────────────────────────────────────────────────────
# 2) 핵심어 추출 + /word_synonyms 호출(문자/JSON 모두 파싱)
# ─────────────────────────────────────────────────────────
_STOP = {
    "the","a","an","and","or","to","of","in","on","for","from","with","that","this",
    "is","are","was","were","be","been","being","by","as","at","it","its","they","their",
    "we","our","you","your","not","but","all","any","can","could","would","should","only",
    "just","more","most","very","also","have","has","had","do","does","did","will","may",
    "into","over","after","before","than","then","there","here","when","which","who","whom",
    "because","so","while"
}
def extract_terms(*texts: str, top: int = 12) -> List[str]:
    cnt = Counter()
    for t in texts:
        for w in re.findall(r"[A-Za-z][A-Za-z'-]{2,}", (t or "").lower()):
            if w in _STOP:
                continue
            cnt[w] += 1
    return [w for w, _ in cnt.most_common(top)]

API_BASE = os.getenv("WORD_API_BASE", "http://127.0.0.1:8000").rstrip("/")

def _parse_bullet_synonym_text(s: str, top_k: int = 3) -> List[dict]:
    lines = [ln.strip() for ln in (s or "").splitlines() if ln.strip()]
    items: List[dict] = []
    cur: dict | None = None

    p_word = re.compile(r"^-+\s*([A-Za-z][A-Za-z\s'-]*)$")
    p_mean = re.compile(r"^-+\s*Meaning\s*:\s*(.+)$", re.I)
    p_syn_head = re.compile(r"^-+\s*Synonyms\s*:\s*$", re.I)
    p_syn_line = re.compile(r"^-+\s*(.+)$")

    for ln in lines:
        m = p_word.match(ln)
        if m and m.group(1).lower() not in ("meaning", "synonyms"):
            if cur and cur.get("word"):
                cur["synonyms"] = cur.get("synonyms", [])[:top_k]
                items.append(cur)
            cur = {"word": m.group(1).strip().lower(), "meaning_ko": "", "synonyms": []}
            continue
        if cur is None:
            continue
        m = p_mean.match(ln)
        if m:
            cur["meaning_ko"] = m.group(1).strip()
            continue
        if p_syn_head.match(ln):
            continue
        m = p_syn_line.match(ln)
        if m:
            raw = m.group(1).strip()
            if not raw:
                continue
            base = raw.split("(")[0].strip().lower()
            have = {syn.split("(")[0].strip().lower() for syn in cur["synonyms"]}
            if base not in have:
                cur["synonyms"].append(raw)
            continue

    if cur and cur.get("word"):
        cur["synonyms"] = cur.get("synonyms", [])[:top_k]
        items.append(cur)
    return items

async def fetch_synonyms_http(terms: List[str], top_k: int = 3) -> List[dict]:
    url = f"{API_BASE}/word_synonyms"
    payload = {"words": terms, "top_k": top_k, "with_meaning": True}
    async with httpx.AsyncClient(timeout=30) as client:
        r = await client.post(url, json=payload)
        r.raise_for_status()
        raw = r.json()

    parsed: List[dict] = []
    if isinstance(raw, dict) and "items" in raw:
        for it in raw["items"]:
            parsed.append({
                "word": (it.get("word") or "").strip().lower(),
                "meaning_ko": (it.get("meaning_ko") or "").strip(),
                "synonyms": [s for s in (it.get("synonyms") or []) if s][:top_k],
            })
    else:
        text_blob = None
        if isinstance(raw, dict):
            for k in ("result","message","text"):
                if isinstance(raw.get(k), str):
                    text_blob = raw[k]; break
            if text_blob is None and len(raw) == 1:
                v = next(iter(raw.values()))  # ← 오타 수정
                if isinstance(v, str): text_blob = v
        elif isinstance(raw, str):
            text_blob = raw
        if not text_blob:
            raise ValueError("Unsupported /word_synonyms response format")
        parsed = _parse_bullet_synonym_text(text_blob, top_k=top_k)

    order = {t.lower(): i for i, t in enumerate(terms)}
    parsed.sort(key=lambda d: order.get(d.get("word", ""), 10**9))
    return parsed

async def get_synonyms_or_fallback(terms: List[str], top_k: int = 3) -> List[dict]:
    try:
        return await fetch_synonyms_http(terms, top_k=top_k)
    except Exception:
        return [{"word": t, "meaning_ko": "", "synonyms": [f"{t}_syn1", f"{t}_syn2", f"{t}_syn3"][:top_k]} for t in terms]

# ─────────────────────────────────────────────────────────
# 3) 구조 분석 호출 (/analyze_structure) → 괄호 적용 텍스트 확보
# ─────────────────────────────────────────────────────────
def _pick_bracketed_text(data: Any) -> str | None:
    if isinstance(data, str):
        return data.strip()
    if isinstance(data, dict):
        for k in ("bracketed","processed_text","result","output","text"):
            v = data.get(k)
            if isinstance(v, str) and v.strip():
                return v.strip()
        for k in ("data","payload","response"):
            v = data.get(k)
            if isinstance(v, dict):
                for kk in ("bracketed","processed_text","result","output","text"):
                    vv = v.get(kk)
                    if isinstance(vv, str) and vv.strip():
                        return vv.strip()
    return None

def _local_bracketize(s: str) -> str:
    if not s:
        return s
    text = s
    # 1) 관계절 [ ... ]
    text = re.sub(
        r"\s+(which|that|who|whom|whose)\b([^,.!?]+)",
        lambda m: " [ " + m.group(1) + m.group(2).rstrip() + " ]",
        text,
        flags=re.IGNORECASE,
    )
    # 2) to부정사 { ... }
    text = re.sub(
        r"\s+to\s+[a-zA-Z]+([^,.!?]*)",
        lambda m: " { to " + m.group(0).strip()[3:] + m.group(1).rstrip() + " }",
        text,
        flags=re.IGNORECASE,
    )
    # 3) 콤마 사이 짧은 전치사구 ( ... )
    def _paren_insertion(m):
        inner = m.group(1).strip()
        if 1 <= len(inner.split()) <= 4:
            return ", ( " + inner + " ) ,"
        return ", " + inner + " ,"
    text = re.sub(
        r",\s*([a-zA-Z]+\s+(?:the|a|an)?\s*[a-zA-Z]+(?:\s+[a-zA-Z]+){0,2})\s*,",
        _paren_insertion,
        text,
    )
    return normalize_bracket_spacing(text)

async def fetch_bracketed_text_http(text: str) -> str:
    url = f"{API_BASE}/analyze_structure"
    payload = {"text": text}
    async with httpx.AsyncClient(timeout=40) as client:
        r = await client.post(url, json=payload)
        r.raise_for_status()
        raw = r.json()
    picked = _pick_bracketed_text(raw)
    if not picked and isinstance(raw, str) and raw.strip():
        picked = raw.strip()
    if not picked:
        raise ValueError("Unsupported /analyze_structure response format")
    return picked

async def get_bracketed_or_fallback(text: str) -> str:
    try:
        return await fetch_bracketed_text_http(text)
    except Exception as e:
        print("[analyze_structure] API failed → fallback to local bracketizer:", repr(e))
        return _local_bracketize(text)

# ─────────────────────────────────────────────────────────
# 4) 요청 모델
# ─────────────────────────────────────────────────────────
class ExportPPTIn(BaseModel):
    passage: str
    passage_bracketed: str | None = None
    date_str: str | None = None
    max_words: int = 12

# ─────────────────────────────────────────────────────────
# 5) PPT 내보내기
# ─────────────────────────────────────────────────────────
@router.post("/ppt")
async def export_ppt(payload: ExportPPTIn):
    try:
        # 5-1) 주제/제목/요지
        tts = analyze_topic_title_summary(payload.passage)

        # 5-2) 구조 분석
        raw = (payload.passage_bracketed or "").strip()
        if raw:
            passage_bracketed = raw
        else:
            passage_bracketed = await get_bracketed_or_fallback(payload.passage)

        # 5-3) 서/본/결
        para = analyze_paragraph(passage_bracketed)
        intro = para.get("outline", {}).get("intro", "")
        body  = para.get("outline", {}).get("body", "")
        concl = para.get("outline", {}).get("conclusion", "")

        # 5-4) 핵심어 & 유의어
        seed_topic = " ".join([tts.get("topic_en",""), tts.get("title_en",""), tts.get("gist_en","")])
        overall_terms = extract_terms(seed_topic, payload.passage, top=payload.max_words)
        intro_terms   = extract_terms(intro, top=4)
        body_terms    = extract_terms(body,  top=4)
        concl_terms   = extract_terms(concl, top=4)

        vocab_overall = await get_synonyms_or_fallback(overall_terms, top_k=3)
        vocab_intro   = await get_synonyms_or_fallback(intro_terms, top_k=3)
        vocab_body    = await get_synonyms_or_fallback(body_terms, top_k=3)
        vocab_concl   = await get_synonyms_or_fallback(concl_terms, top_k=3)

        # 5-5) PPT 작성
        prs = Presentation()

        # Slide 1 — 제목 + 본문 + 요약(ko)
        s1 = prs.slides.add_slide(prs.slide_layouts[6])
        tbox = s1.shapes.add_textbox(Inches(1), Inches(0.8), Inches(9), Inches(1))
        p = tbox.text_frame.paragraphs[0]
        p.text = tts.get("title_en", "")
        p.font.bold = True
        p.font.size = Pt(44)
        p.font.color.rgb = _BLACK

        bodybox = s1.shapes.add_textbox(Inches(1), Inches(2.0), Inches(9), Inches(5))
        tf = bodybox.text_frame; tf.clear(); tf.word_wrap = True
        bp = tf.paragraphs[0]

        clean = normalize_bracket_spacing(passage_bracketed)
        text_for_slide = safe_brackets(clean)
        if not has_any_bracket(text_for_slide) and has_any_bracket(passage_bracketed):
            text_for_slide = safe_brackets(passage_bracketed)

        colorize_brackets(bp, text_for_slide, size_pt=20)

        gist = s1.shapes.add_textbox(Inches(1), Inches(7.2), Inches(9), Inches(0.8))
        gp = gist.text_frame.paragraphs[0]
        gp.text = tts.get("gist_ko","")
        gp.font.size = Pt(18)
        gp.font.color.rgb = _BLACK

        # Slide 2 — 주제/제목/요약 + 서론/본론/결론
        s2 = prs.slides.add_slide(prs.slide_layouts[6])
        badge = s2.shapes.add_textbox(Inches(0.8), Inches(0.7), Inches(3), Inches(0.7))
        bp2 = badge.text_frame.paragraphs[0]
        bp2.text = (payload.date_str or "2025 11 05")
        bp2.font.size = Pt(22)
        bp2.font.color.rgb = _BLACK

        info = s2.shapes.add_textbox(Inches(1), Inches(1.6), Inches(9), Inches(2.6))
        itf = info.text_frame; itf.clear(); itf.word_wrap = True
        def add_pair(label, en, ko):
            a = itf.paragraphs[0] if not itf.paragraphs[0].text else itf.add_paragraph()
            a.text = f"★{label}: {en}"
            a.font.size = Pt(24)
            a.font.color.rgb = _BLACK
            b = itf.add_paragraph()
            b.text = f"({ko})"
            b.font.size = Pt(18)
            b.font.color.rgb = _BLACK
        add_pair("주제", tts.get("topic_en",""), tts.get("topic_ko",""))
        add_pair("제목", tts.get("title_en",""), tts.get("title_ko",""))
        add_pair("요약", tts.get("gist_en",""),  tts.get("gist_ko",""))

        flow = s2.shapes.add_textbox(Inches(1), Inches(4.4), Inches(9), Inches(2.2))
        ft = flow.text_frame; ft.clear(); ft.word_wrap = True
        for label, txt in (("서론", intro), ("본론", body), ("결론", concl)):
            q = ft.paragraphs[0] if not ft.paragraphs[0].text else ft.add_paragraph()
            c = normalize_bracket_spacing(txt)
            t = safe_brackets(c)
            if not has_any_bracket(t) and has_any_bracket(txt):
                t = safe_brackets(txt)
            colorize_brackets(q, f"• {label}: {t}", size_pt=20)

        # Slide 3 — Key Vocabulary & Synonyms
        s3 = prs.slides.add_slide(prs.slide_layouts[1])
        s3.shapes.title.text = "Key Vocabulary & Synonyms"
        tx = s3.shapes.placeholders[1].text_frame; tx.clear(); tx.word_wrap = True

        def add_vocab_item(word: str, syns: list[str]):
            r = tx.paragraphs[0] if not tx.paragraphs[0].text else tx.add_paragraph()
            run = r.add_run()
            run.text = f"{word}"
            run.font.size = Pt(20)
            run.font.bold = True
            run.font.color.rgb = _BLACK
            if syns:
                rr = tx.add_paragraph()
                rr_run = rr.add_run()
                rr_run.text = "  ↳ " + ", ".join(syns)
                rr_run.font.size = Pt(16)
                rr_run.font.color.rgb = _BLACK

        for it in vocab_overall:
            word = (it.get("word") or "").strip()
            syns = [s for s in (it.get("synonyms") or []) if s]
            add_vocab_item(word, syns)

        # Slide 4 — Section-wise Vocabulary
        s4 = prs.slides.add_slide(prs.slide_layouts[1])
        s4.shapes.title.text = "Section-wise Vocabulary"
        bx = s4.shapes.placeholders[1].text_frame; bx.clear(); bx.word_wrap = True

        def add_group(title_txt: str, items: list[dict]):
            head = bx.paragraphs[0] if not bx.paragraphs[0].text else bx.add_paragraph()
            hr = head.add_run()
            hr.text = f"■ {title_txt}"
            hr.font.bold = True
            hr.font.size = Pt(20)
            hr.font.color.rgb = _BLACK
            for it in items:
                word = (it.get("word") or "").strip()
                syns = [s for s in (it.get("synonyms") or []) if s]
                p = bx.add_paragraph()
                r = p.add_run()
                r.text = f"• {word}"
                r.font.size = Pt(18)
                r.font.color.rgb = _BLACK
                if syns:
                    rr = bx.add_paragraph()
                    rr_run = rr.add_run()
                    rr_run.text = "   ↳ " + ", ".join(syns)
                    rr_run.font.size = Pt(16)
                    rr_run.font.color.rgb = _BLACK

        add_group("서론", vocab_intro)
        add_group("본론", vocab_body)
        add_group("결론", vocab_concl)

        bio = BytesIO(); prs.save(bio); bio.seek(0)
        return StreamingResponse(
            bio,
            media_type="application/vnd.openxmlformats-officedocument.presentationml.presentation",
            headers={"Content-Disposition": 'attachment; filename="analysis.pptx"'},
        )
    except Exception as e:
        return JSONResponse(
            status_code=500,
            content={"detail": f"Export PPT failed: {type(e).__name__}: {e}"}
        )